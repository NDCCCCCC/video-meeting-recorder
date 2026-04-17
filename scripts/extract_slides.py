#!/usr/bin/env python3
"""
Extract slide images from PowerPoint files.

This script extracts embedded images from PPTX slides and generates
dual-resolution JPEGs (thumbnails and full-size) for browser preview.

Usage:
    python3 extract_slides.py <pptx_path> <output_dir>

Output:
    JSON: {"success": true, "slide_count": N, "slides": [...]}
    On error: {"success": false, "error": "..."}
"""

import sys
import json
import os
import io
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

# Thumbnail size (16:9 aspect ratio)
THUMBNAIL_SIZE = (200, 112)
# Full-size image resolution (1080p)
FULL_SIZE = (1920, 1080)
# JPEG quality settings
THUMBNAIL_QUALITY = 85
FULL_SIZE_QUALITY = 90


def extract_embedded_images(pptx_path, output_dir):
    """
    Extract embedded images from slides and generate dual-resolution JPEGs.

    Note: This is an MVP implementation that extracts only embedded images.
    Text, shapes, and formatting are NOT rendered (python-pptx limitation).

    Args:
        pptx_path: Path to input PPTX file
        output_dir: Directory to save extracted images

    Returns:
        Tuple (success: bool, result_dict: dict, exit_code: int)
    """
    try:
        # Validate inputs
        if not os.path.exists(pptx_path):
            result = {
                "success": False,
                "error": f"PPTX file not found: {pptx_path}"
            }
            return False, result, 1

        # Create output directories
        thumbnail_dir = os.path.join(output_dir, 'thumbnails')
        fullsize_dir = os.path.join(output_dir, 'fullsize')
        os.makedirs(thumbnail_dir, exist_ok=True)
        os.makedirs(fullsize_dir, exist_ok=True)

        # Load presentation
        prs = Presentation(pptx_path)
        slides_data = []

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_number = slide_idx + 1

            # Look for embedded images in slide shapes
            image_found = False
            for shape in slide.shapes:
                # MSO_SHAPE_TYPE.PICTURE = 13
                if shape.shape_type == 13:
                    try:
                        # Extract image blob
                        image = shape.image
                        image_bytes = image.blob

                        # Open with PIL
                        img = Image.open(io.BytesIO(image_bytes))

                        # Convert to RGB mode (in case it's RGBA or other mode)
                        if img.mode != 'RGB':
                            img = img.convert('RGB')

                        # Generate filenames
                        filename = f'slide_{slide_idx:03d}.jpg'
                        thumbnail_path = os.path.join(thumbnail_dir, filename)
                        fullsize_path = os.path.join(fullsize_dir, filename)

                        # Save full-size image
                        img.save(fullsize_path, 'JPEG', quality=FULL_SIZE_QUALITY)

                        # Save thumbnail
                        img.thumbnail(THUMBNAIL_SIZE, Image.Resampling.LANCZOS)
                        img.save(thumbnail_path, 'JPEG', quality=THUMBNAIL_QUALITY)

                        # Add to results
                        slides_data.append({
                            'slide_number': slide_number,
                            'thumbnail_path': f'thumbnails/{filename}',
                            'fullsize_path': f'fullsize/{filename}'
                        })

                        image_found = True
                        break  # Use first image only
                    except Exception as e:
                        # Skip this shape if extraction fails
                        continue

            # If no image found, create placeholder
            if not image_found:
                filename = f'slide_{slide_idx:03d}.jpg'
                thumbnail_path = os.path.join(thumbnail_dir, filename)
                fullsize_path = os.path.join(fullsize_dir, filename)

                # Create white image with slide number
                img = Image.new('RGB', FULL_SIZE, color='#ffffff')
                draw = ImageDraw.Draw(img)

                # Try to use a default font, fallback to basic text
                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 72)
                except:
                    try:
                        font = ImageFont.truetype("arial.ttf", 72)
                    except:
                        font = ImageFont.load_default()

                # Draw centered text
                text = f"Slide {slide_number}"
                bbox = draw.textbbox((0, 0), text, font=font)
                text_width = bbox[2] - bbox[0]
                text_height = bbox[3] - bbox[1]

                position = ((FULL_SIZE[0] - text_width) // 2, (FULL_SIZE[1] - text_height) // 2)
                draw.text(position, text, fill='black', font=font)

                # Save images
                img.save(fullsize_path, 'JPEG', quality=FULL_SIZE_QUALITY)
                img.thumbnail(THUMBNAIL_SIZE, Image.Resampling.LANCZOS)
                img.save(thumbnail_path, 'JPEG', quality=THUMBNAIL_QUALITY)

                # Add to results
                slides_data.append({
                    'slide_number': slide_number,
                    'thumbnail_path': f'thumbnails/{filename}',
                    'fullsize_path': f'fullsize/{filename}'
                })

        # Return success result
        result = {
            'success': True,
            'slide_count': len(slides_data),
            'slides': slides_data
        }
        return True, result, 0

    except Exception as e:
        result = {
            'success': False,
            'error': str(e)
        }
        return False, result, 1


def main():
    """Main entry point."""
    if len(sys.argv) < 3:
        print(json.dumps({
            'success': False,
            'error': 'Usage: extract_slides.py <pptx_path> <output_dir>'
        }))
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2]

    success, result, exit_code = extract_embedded_images(pptx_path, output_dir)

    # Output result as JSON to stdout
    print(json.dumps(result))
    sys.exit(exit_code)


if __name__ == '__main__':
    main()
