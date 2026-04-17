#!/usr/bin/env python3
"""
Merge slides from multiple PowerPoint files.

This script merges selected slides from multiple PPTX files into a single
output presentation. Slides are copied with their embedded images.

Usage:
    python3 merge_slides.py <output_path> <slide_spec_json>

Where slide_spec_json is a JSON string:
    [{"pptx_path": "...", "slide_numbers": [1,3,5]}, ...]

Output:
    JSON: {"success": true, "slides_merged": N, "output_path": "..."}
    On error: {"success": false, "error": "..."}
"""

import sys
import json
import os
from pptx import Presentation
from pptx.util import Inches

# 16:9 slide dimensions in inches (widescreen)
SLIDE_WIDTH_INCH = 13.333
SLIDE_HEIGHT_INCH = 7.5


def validate_path_safe(path_str):
    """Validate path to prevent traversal attacks.

    Args:
        path_str: Path to validate

    Returns:
        Normalized absolute path

    Raises:
        ValueError: If path contains traversal attempts or is suspicious
    """
    # Normalize the path to resolve any '..' or '.' components
    abs_path = os.path.abspath(path_str)

    # Check if the normalized path is different from what we'd get by resolving components
    # This catches paths like '../../etc/passwd'
    resolved_path = os.path.normpath(path_str)

    # Check for suspicious patterns
    if '..' in path_str.split(os.sep):
        raise ValueError(f"Path contains traversal attempt: {path_str}")

    return abs_path


def merge_slides(output_path, slide_spec):
    """
    Merge slides from multiple PPTX files into a single presentation.

    Args:
        output_path: Path for output merged PPTX file
        slide_spec: List of dicts with pptx_path and slide_numbers

    Returns:
        Tuple (success: bool, result_dict: dict, exit_code: int)
    """
    try:
        # Validate inputs
        if not slide_spec or len(slide_spec) == 0:
            result = {
                'success': False,
                'error': 'No slide specifications provided'
            }
            return False, result, 1

        # Create output presentation with widescreen 16:9 layout
        output_prs = Presentation()
        output_prs.slide_width = Inches(SLIDE_WIDTH_INCH)
        output_prs.slide_height = Inches(SLIDE_HEIGHT_INCH)

        slides_merged = 0

        # Process each source PPTX
        for spec in slide_spec:
            pptx_path = spec.get('pptx_path', '')
            slide_numbers = spec.get('slide_numbers', [])

            # Validate path is safe (no traversal attempts)
            try:
                pptx_path = validate_path_safe(pptx_path)
            except ValueError as e:
                result = {
                    'success': False,
                    'error': str(e)
                }
                return False, result, 1

            # Validate file exists
            if not os.path.exists(pptx_path):
                # Skip silently - continue with other sources
                continue

            try:
                # Load source presentation
                source_prs = Presentation(pptx_path)

                # Process requested slide numbers (1-based)
                for slide_number in slide_numbers:
                    slide_idx = slide_number - 1  # Convert to 0-based

                    # Validate slide index
                    if slide_idx < 0 or slide_idx >= len(source_prs.slides):
                        # Skip invalid slide numbers
                        continue

                    source_slide = source_prs.slides[slide_idx]

                    # Create blank slide in output (layout 6 is blank)
                    blank_layout = output_prs.slide_layouts[6]
                    output_slide = output_prs.slides.add_slide(blank_layout)

                    # Copy picture shapes from source slide
                    for shape in source_slide.shapes:
                        # MSO_SHAPE_TYPE.PICTURE = 13
                        if shape.shape_type == 13:
                            try:
                                # Get image blob
                                img = shape.image
                                img_bytes = img.blob

                                # Add picture to output slide
                                # Position at (0, 0) with full slide size
                                output_slide.shapes.add_picture(
                                    img_bytes,
                                    0,  # left
                                    0,  # top
                                    width=output_prs.slide_width,
                                    height=output_prs.slide_height
                                )
                            except Exception as e:
                                # Skip this shape if copy fails
                                continue

                    slides_merged += 1

            except Exception as e:
                # Skip this source PPTX if loading fails
                continue

        # Ensure at least one slide was merged
        if slides_merged == 0:
            result = {
                'success': False,
                'error': 'No slides were merged from source files'
            }
            return False, result, 1

        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Save output presentation
        output_prs.save(output_path)

        # Return success result
        result = {
            'success': True,
            'slides_merged': slides_merged,
            'output_path': output_path
        }
        return True, result, 0

    except Exception as e:
        result = {
            'success': False,
            'error': str(e)
        }
        return False, result, 1


def main():
    """Main entry point."""
    if len(sys.argv) < 3:
        print(json.dumps({
            'success': False,
            'error': 'Usage: merge_slides.py <output_path> <slide_spec_json>'
        }))
        sys.exit(1)

    output_path = sys.argv[1]
    slide_spec_json = sys.argv[2]

    try:
        slide_spec = json.loads(slide_spec_json)
    except json.JSONDecodeError as e:
        print(json.dumps({
            'success': False,
            'error': f'Invalid JSON in slide_spec: {str(e)}'
        }))
        sys.exit(1)

    success, result, exit_code = merge_slides(output_path, slide_spec)

    # Output result as JSON to stdout
    print(json.dumps(result))
    sys.exit(exit_code)


if __name__ == '__main__':
    main()
