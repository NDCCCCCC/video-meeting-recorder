#!/usr/bin/env python3
"""
Create PowerPoint files from image frames.

This script takes a list of image file paths and creates a PowerPoint
presentation where each image becomes a full-frame slide.

Usage:
    python3 create_pptx.py <output_path> <image1> <image2> ...

Output:
    JSON: {"success": true, "page_count": N, "output_path": "..."}
    On error: {"success": false, "error": "..."}
"""

import sys
import json
import os
from pptx import Presentation
from pptx.util import Inches

# 16:9 slide dimensions in inches
SLIDE_WIDTH_INCH = 10.0
SLIDE_HEIGHT_INCH = 5.625


def create_pptx_from_images(image_paths, output_path):
    """
    Create a PowerPoint file from a list of images.

    Args:
        image_paths: List of image file paths
        output_path: Output PPTX file path

    Returns:
        Tuple (success: bool, result_dict: dict, exit_code: int)
    """
    try:
        # Validate inputs
        if not image_paths:
            result = {
                "success": False,
                "error": "No image paths provided"
            }
            return False, result, 1

        # Create presentation with 16:9 slide size
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_WIDTH_INCH)
        prs.slide_height = Inches(SLIDE_HEIGHT_INCH)

        page_count = 0
        skipped = []
        missing_files = []

        # Add each image as a slide
        for i, img_path in enumerate(image_paths):
            try:
                # Check if file exists
                if not os.path.exists(img_path):
                    missing_files.append((i, img_path))
                    skipped.append(img_path)
                    continue

                # Add blank slide (layout 6 is the blank layout)
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)

                # Add image to slide (full-frame, no margins)
                # Position at (0, 0) and size to full slide dimensions
                slide.shapes.add_picture(
                    img_path,
                    0,  # left
                    0,  # top
                    width=prs.slide_width,
                    height=prs.slide_height
                )

                page_count += 1

            except Exception as e:
                # Skip silently - don't output errors to stdout
                skipped.append(img_path)
                continue

        # Ensure at least one slide was created
        if page_count == 0:
            result = {
                "success": False,
                "error": f"No valid slides created from {len(image_paths)} input images. {len(missing_files)} files were missing.",
                "skipped": skipped,
                "missing_files": missing_files,
                "input_count": len(image_paths)
            }
            return False, result, 1

        # Warn if significant number of files were missing
        if len(missing_files) > 0:
            # Use stderr for warnings so stdout remains pure JSON
            import sys
            print(f"Warning: {len(missing_files)} out of {len(image_paths)} image files were missing", file=sys.stderr)
            print(f"Missing files: {[p for _, p in missing_files[:5]]}{'...' if len(missing_files) > 5 else ''}", file=sys.stderr)

        # Ensure output directory exists
        output_dir = os.path.dirname(output_path)
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)

        # Save presentation
        prs.save(output_path)

        # Return success result
        result = {
            "success": True,
            "page_count": page_count,
            "output_path": output_path,
            "skipped_count": len(skipped),
            "input_count": len(image_paths),
            "missing_files": missing_files if len(missing_files) > 0 else None
        }
        return True, result, 0

    except Exception as e:
        # Return error result with more details
        import traceback
        result = {
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }
        return False, result, 1


def main():
    """Main entry point."""
    if len(sys.argv) < 3:
        print(json.dumps({
            "success": False,
            "error": "Usage: create_pptx.py <output_path> <image1> <image2> ..."
        }), file=sys.stderr)
        sys.exit(1)

    output_path = sys.argv[1]
    image_paths = sys.argv[2:]

    success, result, exit_code = create_pptx_from_images(image_paths, output_path)

    # Output result as JSON
    if success:
        print(json.dumps(result))
    else:
        print(json.dumps(result), file=sys.stderr)

    sys.exit(exit_code)


if __name__ == "__main__":
    main()
